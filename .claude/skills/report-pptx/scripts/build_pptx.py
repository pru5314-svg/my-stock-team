# -*- coding: utf-8 -*-
"""
report-pptx : 종목 리서치 마크다운 → 투자증권 리서치센터 톤 PPTX 변환기.

입력 : reports/{종목명}.md
출력 : reports/pptx/{종목명}.pptx  (pptx 폴더 없으면 생성)

슬라이드 순서(고정): 표지 → 종목 개요 → 재무 요약 → 가격/추세 → 뉴스·심리 → 리스크 → 한 줄 종합
디자인           : 화이트 배경 / 딥네이비 제목 / 그레이 본문 / 한화 오렌지 절제 강조 / 맑은 고딕 고정

사용:
    python build_pptx.py "삼성전자"
    python build_pptx.py --input reports/삼성전자.md --output reports/pptx/삼성전자.pptx --date 2026-06-16
"""
import argparse
import datetime
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
except ImportError:
    sys.stderr.write(
        "[report-pptx] python-pptx 가 필요합니다. 설치: pip install python-pptx\n"
    )
    sys.exit(2)

# ---------------------------------------------------------------- 디자인 상수
FONT = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x2A, 0x44)      # 제목/표 헤더
ORANGE = RGBColor(0xF3, 0x73, 0x21)    # 상단 바·구분선·키 수치 강조 (절제)
BODY = RGBColor(0x40, 0x40, 0x40)      # 본문 그레이
MUTED = RGBColor(0x8C, 0x8C, 0x8C)     # 출처/캡션 등 약한 그레이
RULE = "D9D9D9"                        # 표 본문 괘선 (얇은 그레이)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
BODY_TOP = Inches(1.35)                # 제목/구분선 아래 본문 시작
BODY_BOTTOM = Inches(7.0)             # 본문 하단 한계(이 아래로 표가 넘치면 안 됨)

DISCLAIMER = "본 자료는 정보 제공 목적이며, 투자 판단의 책임은 투자자 본인에게 있습니다."

# 매수/매도·목표가 단정 금지 — 린트 대상
FORBIDDEN = [
    "매수", "매도", "목표가", "강력매수", "강력 매수", "적극매수", "적극 매수",
    "비중확대", "비중축소", "비중 확대", "비중 축소", "BUY", "SELL", "target price",
]

# 고정 슬라이드 정의: (제목, 매핑 키워드들)
SECTIONS = [
    ("종목 개요", ["개요", "기업", "소개", "overview"]),
    ("재무 요약", ["재무", "실적", "financ"]),
    ("가격/추세", ["가격", "추세", "주가", "기술", "차트", "price"]),
    ("뉴스·심리", ["뉴스", "심리", "센티", "news", "sentiment"]),
    ("리스크", ["리스크", "위험", "risk"]),
    ("한 줄 종합", ["종합", "결론", "한 줄", "한줄", "conclusion"]),
]


# ---------------------------------------------------------------- 폰트 헬퍼
def style_run(run, size=None, color=None, bold=None, name=FONT):
    """latin/ea/cs 모두 맑은 고딕으로 강제 → 한글 깨짐 방지."""
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


# 키 수치 강조: **...** → 오렌지 볼드
EMPH = re.compile(r"\*\*(.+?)\*\*")


def add_rich_text(paragraph, text, size, color, bold=False):
    """문장 안의 **강조**를 오렌지 볼드로 분할 렌더."""
    pos = 0
    found = False
    for m in EMPH.finditer(text):
        if m.start() > pos:
            r = paragraph.add_run(); r.text = text[pos:m.start()]
            style_run(r, size=size, color=color, bold=bold)
        r = paragraph.add_run(); r.text = m.group(1)
        style_run(r, size=size, color=ORANGE, bold=True)
        pos = m.end()
        found = True
    if pos < len(text):
        r = paragraph.add_run(); r.text = text[pos:]
        style_run(r, size=size, color=color, bold=bold)
    if not found and pos == 0:  # 빈 문자열 방지
        r = paragraph.add_run(); r.text = text
        style_run(r, size=size, color=color, bold=bold)


def strip_emph(text):
    return EMPH.sub(r"\1", text)


# ---------------------------------------------------------------- 셀 테두리
def set_cell_border_bottom(cell, color=RULE, width_emu=9525):
    """본문 행 하단에만 얇은 그레이 괘선."""
    tcPr = cell._tc.get_or_add_tcPr()
    for e in tcPr.findall(qn("a:lnB")):
        tcPr.remove(e)
    ln = tcPr.makeelement(qn("a:lnB"),
                          {"w": str(width_emu), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
    fill.append(clr)
    ln.append(fill)
    tcPr.insert(0, ln)  # 라인은 fill 류보다 앞에 와야 함


# ---------------------------------------------------------------- 마크다운 파싱
def is_sep_row(cells):
    return all(re.fullmatch(r":?-{2,}:?", c.strip() or "-") for c in cells) and cells


def parse_table(lines):
    rows = []
    for ln in lines:
        s = ln.strip()
        if s.startswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            if is_sep_row(cells):
                continue
            rows.append(cells)
    # 열 수 정규화
    if rows:
        width = max(len(r) for r in rows)
        rows = [r + [""] * (width - len(r)) for r in rows]
    return rows


IMG = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")


def parse_blocks(lines):
    """섹션 본문을 등장 순서대로 (type, payload) 블록 리스트로 분해."""
    blocks = []
    i = 0
    n = len(lines)
    while i < n:
        s = lines[i].strip()
        if not s:
            i += 1
            continue
        m = IMG.search(s)
        if m:
            blocks.append(("image", m.group(1)))
            i += 1
            continue
        if s.startswith("|"):
            j = i
            tbl = []
            while j < n and lines[j].strip().startswith("|"):
                tbl.append(lines[j])
                j += 1
            rows = parse_table(tbl)
            if rows:
                blocks.append(("table", rows))
            i = j
            continue
        # 텍스트 묶음 (다음 표/이미지/빈줄 전까지)
        txt = []
        while i < n:
            t = lines[i].strip()
            if not t or t.startswith("|") or IMG.search(t):
                break
            txt.append(t)
            i += 1
        if txt:
            blocks.append(("text", txt))
    return blocks


def parse_markdown(md):
    lines = md.splitlines()
    title = None
    date = None
    # 제목(첫 H1) + 작성일
    for ln in lines:
        s = ln.strip()
        if title is None and s.startswith("# "):
            title = s[2:].strip()
        if date is None:
            m = re.search(r"(작성일|기준일)\s*[:：]?\s*(\d{4}[-./]\d{1,2}[-./]\d{1,2})", s)
            if m:
                date = m.group(2).replace(".", "-").replace("/", "-")
    # 섹션 분해 (## 기준)
    sections = {}
    cur = None
    buf = []
    for ln in lines:
        s = ln.strip()
        if s.startswith("## "):
            if cur is not None:
                sections[cur] = buf
            cur = s[3:].strip()
            buf = []
        elif cur is not None:
            buf.append(ln)
    if cur is not None:
        sections[cur] = buf
    return title, date, sections


def match_section(sections, keywords):
    for name, body in sections.items():
        low = name.lower()
        if any(k.lower() in low for k in keywords):
            return body
    return None


# ---------------------------------------------------------------- 슬라이드 빌드
def blank_slide(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def add_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def build_cover(prs, title, date):
    slide = blank_slide(prs)
    # 상단 오렌지 바
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.3), ORANGE)
    # 종목명
    _, tf = add_textbox(slide, MARGIN, Inches(2.7), SLIDE_W - 2 * MARGIN, Inches(1.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    style_run(r, size=44, color=NAVY, bold=True)
    # 부제
    _, tf2 = add_textbox(slide, MARGIN, Inches(4.1), SLIDE_W - 2 * MARGIN, Inches(0.6))
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = "종목 리서치 리포트  |  작성일 %s" % date
    style_run(r2, size=16, color=BODY)
    # 오렌지 짧은 구분선
    add_rect(slide, MARGIN, Inches(2.55), Inches(1.2), Inches(0.06), ORANGE)
    # 하단 면책
    _, tf3 = add_textbox(slide, MARGIN, Inches(6.9), SLIDE_W - 2 * MARGIN, Inches(0.4))
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run(); r3.text = DISCLAIMER
    style_run(r3, size=10, color=MUTED)
    return slide


def add_title_and_rule(slide, title):
    _, tf = add_textbox(slide, MARGIN, Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.7))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    style_run(r, size=26, color=NAVY, bold=True)
    # 오렌지 구분선
    add_rect(slide, MARGIN, Inches(1.18), SLIDE_W - 2 * MARGIN, Inches(0.045), ORANGE)


def render_text_block(slide, y, lines):
    """불릿/문단 텍스트 렌더. 차지한 높이(Emu) 반환."""
    height = Inches(min(5.0, 0.32 * len(lines) + 0.2))
    _, tf = add_textbox(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, height)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        bullet = ""
        body = ln
        m = re.match(r"^(\d+\.|[-*•])\s+(.*)$", ln)
        if m:
            bullet = "• "
            body = m.group(2)
        # 출처 표기는 약한 그레이로 보이도록 그대로 두되 본문 색 사용
        if bullet:
            rb = p.add_run(); rb.text = bullet
            style_run(rb, size=14, color=ORANGE, bold=True)
        add_rich_text(p, body, size=14, color=BODY)
    return height


def _max_body_height():
    return BODY_BOTTOM - BODY_TOP


def render_table_block(slide, y, rows):
    """
    표 렌더. 슬라이드 밖으로 잘리지 않도록:
      1) 가용 높이에 맞춰 글꼴/행높이를 단계적으로 축소
      2) 그래도 넘치면 행을 잘라내고 마지막에 '이하 N행 생략' 안내행 추가
    차지한 높이(Emu) 반환.
    """
    avail = (BODY_BOTTOM - y)
    n_cols = max(len(r) for r in rows)
    header, body_rows = rows[0], rows[1:]

    # 폰트/행높이 후보 (큰 것부터). (본문pt, 행높이in)
    candidates = [(13, 0.42), (12, 0.36), (11, 0.32), (10, 0.28), (9, 0.25), (8, 0.23)]
    header_h = 0.0
    chosen = candidates[-1]
    truncated = 0
    kept = body_rows
    for fsize, rh in candidates:
        head_rh = rh + 0.04
        total = Emu(int((head_rh + rh * len(body_rows)) * Inches(1)))
        if total <= avail:
            chosen = (fsize, rh)
            header_h = head_rh
            kept = body_rows
            truncated = 0
            break
    else:
        # 가장 작은 글꼴로도 안 들어가면 → 행 잘라내기
        fsize, rh = candidates[-1]
        head_rh = rh + 0.04
        avail_in = avail / Inches(1)
        max_body = int((avail_in - head_rh) / rh) - 1  # 생략 안내행 1줄 확보
        max_body = max(max_body, 1)
        if max_body < len(body_rows):
            truncated = len(body_rows) - max_body
            kept = body_rows[:max_body]
        chosen = (fsize, rh)
        header_h = head_rh

    fsize, rh = chosen
    display = [header] + kept
    if truncated:
        note = ["이하 %d개 행 생략 (전체는 원본 마크다운 참조)" % truncated] + [""] * (n_cols - 1)
        display.append(note)

    n_rows = len(display)
    table_w = SLIDE_W - 2 * MARGIN
    table_h = Emu(int((header_h + rh * (n_rows - 1)) * Inches(1)))
    gf = slide.shapes.add_table(n_rows, n_cols, MARGIN, y, table_w, table_h)
    tbl = gf.table

    # 기본 테마 스타일/밴딩 제거 → 우리가 직접 색칠
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)

    # 행 높이
    tbl.rows[0].height = Inches(header_h)
    for ri in range(1, n_rows):
        tbl.rows[ri].height = Inches(rh)

    for ri in range(n_rows):
        is_header = (ri == 0)
        is_note = truncated and ri == n_rows - 1
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt = display[ri][ci] if ci < len(display[ri]) else ""
            # 채움
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_header else WHITE
            if not is_header:
                set_cell_border_bottom(cell)
            # 텍스트
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            if is_header:
                r = p.add_run(); r.text = strip_emph(txt)
                style_run(r, size=fsize, color=WHITE, bold=True)
            elif is_note:
                r = p.add_run(); r.text = txt
                style_run(r, size=max(fsize - 1, 8), color=MUTED, bold=False)
                p.alignment = PP_ALIGN.LEFT
            else:
                add_rich_text(p, txt, size=fsize, color=BODY)
    return table_h


def render_image_block(slide, y, path, md_dir, root):
    """이미지 임베드. 가용 영역에 맞춰 비율 유지 축소. 차지 높이 반환."""
    candidates = [path,
                  os.path.join(md_dir, path),
                  os.path.join(root, path)]
    found = next((c for c in candidates if os.path.isfile(c)), None)
    avail_h = BODY_BOTTOM - y
    if not found:
        _, tf = add_textbox(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.5))
        r = tf.paragraphs[0].add_run()
        r.text = "[차트 이미지를 찾을 수 없음: %s]" % path
        style_run(r, size=12, color=MUTED)
        return Inches(0.5)
    max_w = SLIDE_W - 2 * MARGIN
    max_h = min(avail_h, Inches(4.6))
    pic = slide.shapes.add_picture(found, MARGIN, y, width=max_w)
    if pic.height > max_h:  # 높이 초과 시 비율 유지 재스케일
        ratio = max_h / pic.height
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        pic.left = int((SLIDE_W - pic.width) / 2)
    return Emu(int(pic.height)) + Inches(0.1)


def build_content(prs, title, blocks, md_dir, root):
    slide = blank_slide(prs)
    add_title_and_rule(slide, title)
    y = BODY_TOP
    if not blocks:
        _, tf = add_textbox(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.6))
        r = tf.paragraphs[0].add_run()
        r.text = "내용 없음 — 입력 마크다운에 해당 섹션이 없습니다."
        style_run(r, size=14, color=MUTED)
        return slide
    for btype, payload in blocks:
        if y >= BODY_BOTTOM - Inches(0.4):
            break  # 더 넣을 공간 없음 (안전 차단)
        if btype == "text":
            used = render_text_block(slide, y, payload)
        elif btype == "table":
            used = render_table_block(slide, y, payload)
        elif btype == "image":
            used = render_image_block(slide, y, payload, md_dir, root)
        else:
            used = Inches(0)
        y = Emu(int(y) + int(used) + int(Inches(0.18)))
    return slide


# ---------------------------------------------------------------- 린트
def lint(md, title):
    warnings = []
    low = md.lower()
    for term in FORBIDDEN:
        if term.lower() in low:
            warnings.append("금지 표현 발견: '%s' — 매수/매도·목표가 단정은 제거하세요." % term)
    # 출처 없는 수치(아주 가벼운 휴리스틱): 숫자가 있는 표/문단에 '출처'가 전혀 없으면 경고
    if re.search(r"\d", md) and "출처" not in md and "기준일" not in md:
        warnings.append("수치는 있으나 '출처/기준일' 표기가 보이지 않습니다 — 출처 없는 숫자는 싣지 마세요.")
    return warnings


# ---------------------------------------------------------------- 메인
def main():
    ap = argparse.ArgumentParser(description="리서치 마크다운 → 투자증권 톤 PPTX")
    ap.add_argument("name", nargs="?", help="종목명 (reports/{종목명}.md)")
    ap.add_argument("--input", help="입력 마크다운 경로")
    ap.add_argument("--output", help="출력 PPTX 경로")
    ap.add_argument("--date", help="작성일 YYYY-MM-DD (마크다운에 없을 때)")
    ap.add_argument("--root", default=os.getcwd(), help="프로젝트 루트 (기본: 현재 디렉터리)")
    args = ap.parse_args()

    root = args.root
    if args.input:
        in_path = args.input
        name = args.name or os.path.splitext(os.path.basename(in_path))[0]
    elif args.name:
        name = args.name
        in_path = os.path.join(root, "reports", "%s.md" % name)
    else:
        ap.error("종목명 또는 --input 중 하나는 필요합니다.")

    if not os.path.isfile(in_path):
        sys.stderr.write("[report-pptx] 입력 파일이 없습니다: %s\n" % in_path)
        sys.exit(1)

    with open(in_path, encoding="utf-8") as f:
        md = f.read()

    title, date, sections = parse_markdown(md)
    title = title or name
    date = args.date or date or datetime.date.today().isoformat()

    out_path = args.output or os.path.join(root, "reports", "pptx", "%s.pptx" % name)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    md_dir = os.path.dirname(os.path.abspath(in_path))

    # 린트
    for w in lint(md, title):
        sys.stderr.write("[경고] %s\n" % w)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs, title, date)
    for sec_title, keywords in SECTIONS:
        body = match_section(sections, keywords)
        blocks = parse_blocks(body) if body else []
        build_content(prs, sec_title, blocks, md_dir, root)

    prs.save(out_path)
    print("[report-pptx] 생성 완료 → %s (슬라이드 %d장)" % (out_path, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
